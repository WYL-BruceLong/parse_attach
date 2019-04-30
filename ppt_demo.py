# #!/usr/bin/python3
# -*- coding: utf-8 -*-
# @Time : 2019-04-30 22:09
# @Author : BruceLong
# @FileName: ppt_demo.py
# @Email   : 18656170559@163.com
# @Software: PyCharm
# @Blog ：http://www.cnblogs.com/yunlongaimeng/
from pptx import Presentation


def read_pptx(path):
    prs = Presentation(path)
    MBslides = prs.slides
    MBslidesNum = len(MBslides)
    data = ''
    for i in range(0, MBslidesNum):
        MBshape = MBslides[i].shapes
        MBshapeNum = len(MBshape)
        for j in range(0, MBshapeNum):
            if MBshape[j].has_text_frame:
                frame_text = MBshape[j].text
                data += frame_text.strip().replace("\n", '')
            elif MBshape[j].has_table:
                for cell in MBshape[j].table.iter_cells():
                    table_text = cell.text.strip().replace("\n", '')
                    data += table_text
    print(data)


if __name__ == '__main__':
    # path = r'C:\Users\hbsxw\Desktop\git_projects\parse_attach\attach_files\pptx_7107817260a1894b7ee0b468b54e29e9_日本站月度大促3月到6月安排.pptx'
    # path = r'C:\Users\hbsxw\Desktop\git_projects\parse_attach\attach_files\PPTX_68d034dc46e22b8bb719dbcb77370f8a_Coupon优惠券使用说明.PPTX'
    path = r'C:\Users\hbsxw\Desktop\git_projects\parse_attach\attach_files\pptx_7107817260a1894b7ee0b468b54e29e9_日本站月度大促3月到6月安排.ppt'
    read_pptx(path)
